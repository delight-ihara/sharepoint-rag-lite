"""SharePoint → テキスト抽出 → pgvector インジェストパイプライン

既存 sp_to_blob.py の ACL ロジックを流用。
Blob Storage を経由せず、直接 pgvector に書き込む。
"""

import hashlib
import io
import logging
import sys
import time
from urllib.parse import quote

import chardet
import fitz  # PyMuPDF
import requests
from docx import Document as DocxDocument
from openai import APIConnectionError, APIStatusError, AzureOpenAI, RateLimitError
from openpyxl import load_workbook
from tenacity import (
    before_sleep_log,
    retry,
    retry_if_exception,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
)

from .acl import check_graph_permissions, get_app_token, resolve_folder_acl
from .config import (
    AZURE_OPENAI_ENDPOINT,
    AZURE_OPENAI_KEY,
    EMBEDDING_DEPLOYMENT,
    GRAPH_BASE,
    MAX_FILE_SIZE_MB,
    SP_DRIVE_ID,
    TARGET_FOLDERS,
)
from .db import get_conn, put_conn

log = logging.getLogger(__name__)


def _is_transient_openai_error(exc: BaseException) -> bool:
    """リトライ対象の OpenAI 一時エラーか判定"""
    if isinstance(exc, (RateLimitError, APIConnectionError)):
        return True
    if isinstance(exc, APIStatusError) and exc.status_code in (502, 503, 504):
        return True
    return False


def _is_transient_http_error(exc: BaseException) -> bool:
    """リトライ対象の HTTP 一時エラーか判定"""
    if isinstance(exc, (requests.exceptions.ConnectionError, requests.exceptions.Timeout)):
        return True
    if isinstance(exc, requests.exceptions.HTTPError) and exc.response is not None:
        if exc.response.status_code in (429, 502, 503, 504):
            return True
    return False


# get_app_token() は src/acl.py に移動済み


# resolve_folder_acl() は src/acl.py に移動済み


# ── SP ファイル取得 ─────────────────────────────────────

def list_sp_files(token: str) -> list[dict]:
    """SharePoint ドキュメントライブラリの全ファイルを再帰取得（ページネーション対応）"""
    headers = {"Authorization": f"Bearer {token}"}
    files = []
    max_size = MAX_FILE_SIZE_MB * 1024 * 1024

    def walk(path: str = ""):
        if path:
            url: str | None = f"{GRAPH_BASE}/drives/{SP_DRIVE_ID}/root:/{quote(path)}:/children?$top=200"
        else:
            url = f"{GRAPH_BASE}/drives/{SP_DRIVE_ID}/root/children?$top=200"

        while url:
            resp = requests.get(url, headers=headers, timeout=30)
            resp.raise_for_status()
            data = resp.json()

            for item in data.get("value", []):
                if item.get("folder"):
                    child_path = f"{path}/{item['name']}" if path else item["name"]
                    if not path and TARGET_FOLDERS:
                        if not any(child_path.startswith(prefix) for prefix in TARGET_FOLDERS):
                            continue
                    walk(child_path)
                elif item.get("file"):
                    size = item.get("size", 0)
                    if size > max_size:
                        log.warning("  ファイルサイズ超過 (%dMB > %dMB): %s", size // (1024*1024), MAX_FILE_SIZE_MB, item["name"])
                        continue
                    files.append({
                        "id": item["id"],
                        "name": item["name"],
                        "path": path,
                        "size": size,
                        "mimeType": item["file"].get("mimeType", ""),
                        "lastModified": item.get("lastModifiedDateTime", ""),
                        "webUrl": item.get("webUrl", ""),
                        "@microsoft.graph.downloadUrl": item.get("@microsoft.graph.downloadUrl", ""),
                    })

            url = data.get("@odata.nextLink")

    walk()
    return files


@retry(
    retry=retry_if_exception(_is_transient_http_error),
    wait=wait_exponential(multiplier=1, min=1, max=10),
    stop=stop_after_attempt(3),
    before_sleep=before_sleep_log(log, logging.WARNING),
    reraise=True,
)
def download_file(token: str, file_id: str) -> bytes:
    """ファイルのバイナリをダウンロード（リトライ付き）"""
    headers = {"Authorization": f"Bearer {token}"}
    url = f"{GRAPH_BASE}/drives/{SP_DRIVE_ID}/items/{file_id}/content"
    resp = requests.get(url, headers=headers, timeout=60)
    resp.raise_for_status()
    return resp.content


# ── テキスト抽出 ───────────────────────────────────────

def extract_text(content: bytes, filename: str) -> str:
    """ファイルからテキストを抽出"""
    # ~$ で始まるファイルはOfficeの一時ファイル → スキップ
    if filename.startswith("~$"):
        log.info("  一時ファイルをスキップ: %s", filename)
        return ""

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        if ext == "pdf":
            return _extract_pdf(content)
        elif ext == "docx":
            return _extract_docx(content)
        elif ext == "xlsx":
            return _extract_xlsx(content)
        elif ext == "pptx":
            return _extract_pptx(content)
        elif ext in ("txt", "csv", "md", "json", "xml", "html", "htm"):
            return _decode_text(content, filename)
        elif ext == "xls":
            return _extract_xls(content)
        elif ext in ("doc", "ppt"):
            return _extract_legacy_office(content, filename, ext)
        else:
            log.warning("  未対応形式 (.%s): %s", ext, filename)
            return ""
    except Exception as e:
        err_str = str(e).lower()
        if any(kw in err_str for kw in ("password", "encrypt", "protected", "decrypt")):
            log.warning("  パスワード保護ファイル — テキスト抽出不可: %s", filename)
        else:
            log.warning("  テキスト抽出エラー: %s — %s", filename, e)
        return ""


def _extract_xls(content: bytes) -> str:
    """旧 .xls 形式を xlrd でテキスト抽出（クロスプラットフォーム）"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=content)
        texts = []
        for sheet in wb.sheets():
            texts.append(f"=== {sheet.name} ===")
            for row_idx in range(sheet.nrows):
                vals = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                if any(v.strip() for v in vals):
                    texts.append("\t".join(vals))
        return "\n".join(texts)
    except ImportError:
        log.warning("  xlrd 未インストール — .xls テキスト抽出不可")
        return _extract_legacy_office(content, "file.xls", "xls")
    except Exception as e:
        log.warning("  .xls 抽出エラー: %s — フォールバック試行", e)
        return _extract_legacy_office(content, "file.xls", "xls")


def _extract_legacy_office(content: bytes, filename: str, ext: str) -> str:
    """旧Office形式 (.doc/.xls/.ppt) の汎用テキスト抽出

    1. win32com (Windows) — Word/Excel/PowerPoint COM
    2. LibreOffice headless (Linux/Docker) — soffice --convert-to txt
    """
    import tempfile
    import os

    # === 1. win32com (Windows) ===
    tmp_path = None
    try:
        import win32com.client
        import pythoncom

        pythoncom.CoInitialize()

        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
            f.write(content)
            tmp_path = f.name

        app_map = {
            "doc": ("Word.Application", "Documents", lambda app, path: _com_extract_word(app, path)),
            "xls": ("Excel.Application", "Workbooks", lambda app, path: _com_extract_excel(app, path)),
            "ppt": ("PowerPoint.Application", "Presentations", lambda app, path: _com_extract_ppt(app, path)),
        }

        if ext not in app_map:
            return ""

        app_name, _, extractor = app_map[ext]
        app = win32com.client.Dispatch(app_name)
        app.Visible = False
        try:
            return extractor(app, tmp_path)
        finally:
            app.Quit()
            pythoncom.CoUninitialize()

    except ImportError:
        # win32com なし → LibreOffice フォールバック
        pass
    except Exception as e:
        log.warning("  COM抽出エラー (.%s): %s — LibreOffice フォールバック試行", ext, e)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

    # === 2. LibreOffice headless (Linux/Docker) ===
    return _extract_with_libreoffice(content, filename, ext)


def _com_extract_word(app, path: str) -> str:
    doc = app.Documents.Open(path, ReadOnly=True)
    text = doc.Content.Text
    doc.Close(False)
    return text


def _com_extract_excel(app, path: str) -> str:
    wb = app.Workbooks.Open(path, ReadOnly=True)
    texts = []
    for ws in wb.Worksheets:
        texts.append(f"=== {ws.Name} ===")
        used = ws.UsedRange
        if used:
            for row in used.Rows:
                vals = [str(cell.Value) if cell.Value is not None else "" for cell in row.Cells]
                if any(v.strip() for v in vals):
                    texts.append("\t".join(vals))
    wb.Close(False)
    return "\n".join(texts)


def _com_extract_ppt(app, path: str) -> str:
    prs = app.Presentations.Open(path, ReadOnly=True, WithWindow=False)
    texts = []
    for slide in prs.Slides:
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                texts.append(shape.TextFrame.TextRange.Text)
    prs.Close()
    return "\n".join(texts)


def _extract_with_libreoffice(content: bytes, filename: str, ext: str) -> str:
    """LibreOffice headless でテキスト変換（Linux/Docker 用）"""
    import tempfile
    import os
    import subprocess
    import shutil

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        log.warning("  win32com も LibreOffice も利用不可 — .%s テキスト抽出不可: %s", ext, filename)
        return ""

    tmp_dir = tempfile.mkdtemp()
    tmp_path = os.path.join(tmp_dir, f"input.{ext}")
    try:
        with open(tmp_path, "wb") as f:
            f.write(content)

        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", tmp_dir, tmp_path],
            capture_output=True, timeout=60,
        )
        txt_path = os.path.join(tmp_dir, "input.txt")
        if os.path.exists(txt_path):
            with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        log.warning("  LibreOffice 変換失敗: %s (rc=%d)", filename, result.returncode)
        return ""
    except Exception as e:
        log.warning("  LibreOffice 抽出エラー: %s — %s", filename, e)
        return ""
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def _decode_text(content: bytes, filename: str) -> str:
    """テキストファイルのエンコーディングを自動検出してデコード"""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        detected = chardet.detect(content)
        encoding = detected.get("encoding") or "utf-8"
        log.info("  エンコーディング自動検出: %s → %s (confidence=%.0f%%)",
                 filename, encoding, (detected.get("confidence") or 0) * 100)
        return content.decode(encoding, errors="replace")


def _clean_pdf_text(text: str) -> str:
    """PDF 抽出テキストからノイズを除去

    - ドットリーダー行（目次の「第1条 目的 . . . . . . 3」等）を除去
    - 連続ドットを圧縮
    - 意味のある文字がほぼ無い行を除去
    """
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        # ドットリーダー行の検出: ドットが5個以上連続（間にスペース可）
        stripped = line.strip()
        if re.search(r'(?:\.\s*){5,}', stripped):
            # 意味のある文字（ドット・スペース・数字以外）が残るか確認
            meaningful = re.sub(r'[\.\s\d\u3000]+', '', stripped)
            if len(meaningful) < 3:
                # ほぼドットのみ → 完全除去
                continue
            else:
                # 目次行: ドットリーダー部分だけ除去して見出しテキストを保持
                cleaned_line = re.sub(r'\s*(?:\.\s*){3,}\s*\d*\s*$', '', stripped)
                cleaned_line = cleaned_line.rstrip('. \t')
                if cleaned_line.strip():
                    cleaned.append(cleaned_line.strip())
                continue
        # 行がドットとスペースだけ → 除去
        if stripped and not re.sub(r'[\.\s\d\u3000]+', '', stripped):
            continue
        cleaned.append(line)
    return "\n".join(cleaned)


def _extract_pdf(content: bytes) -> str:
    doc = fitz.open(stream=content, filetype="pdf")
    if doc.is_encrypted:
        log.warning("  パスワード保護PDF — テキスト抽出不可")
        doc.close()
        return ""
    texts = [page.get_text() for page in doc]
    doc.close()
    raw = "\n".join(texts)
    cleaned = _clean_pdf_text(raw)

    # テキストが空 or 極めて少ない → 画像PDF の可能性 → OCR フォールバック
    if len(cleaned.strip()) < 50:
        ocr_text = _ocr_pdf(content)
        if ocr_text and len(ocr_text.strip()) > len(cleaned.strip()):
            log.info("  OCR フォールバック使用（画像PDF）")
            return ocr_text

    return cleaned


def _ocr_pdf(content: bytes) -> str:
    """Tesseract OCR で画像 PDF からテキスト抽出"""
    import os
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        log.warning("  pytesseract/Pillow 未インストール — OCR 不可")
        return ""

    # Tesseract パス設定（Windows）
    tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path

    # TESSDATA_PREFIX が未設定ならユーザーディレクトリを試行
    if not os.environ.get("TESSDATA_PREFIX"):
        fallback = os.path.expanduser("~/tessdata")
        if os.path.isdir(fallback):
            os.environ["TESSDATA_PREFIX"] = fallback

    try:
        doc = fitz.open(stream=content, filetype="pdf")
        texts = []
        for page in doc:
            # ページを画像に変換（300 DPI）
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            text = pytesseract.image_to_string(img, lang="jpn+eng")
            texts.append(text)
        doc.close()
        return "\n".join(texts)
    except Exception as e:
        log.warning("  OCR 失敗: %s", e)
        return ""


def _extract_docx(content: bytes) -> str:
    try:
        doc = DocxDocument(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except KeyError:
        # comments.xml 等が欠損した壊れた docx → zipfile + XML で直接抽出
        return _extract_docx_raw(content)


def _extract_docx_raw(content: bytes) -> str:
    """壊れた .docx から zipfile + XML で本文を直接抽出"""
    import zipfile
    from xml.etree import ElementTree

    try:
        with zipfile.ZipFile(io.BytesIO(content)) as z:
            if "word/document.xml" not in z.namelist():
                return ""
            xml_content = z.read("word/document.xml")
            tree = ElementTree.fromstring(xml_content)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            paragraphs = []
            for p in tree.iter(f"{{{ns['w']}}}p"):
                texts = [t.text for t in p.iter(f"{{{ns['w']}}}t") if t.text]
                if texts:
                    paragraphs.append("".join(texts))
            return "\n".join(paragraphs)
    except Exception as e:
        log.warning("  docx raw抽出も失敗: %s", e)
        return ""


def _extract_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        texts.append(f"=== {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            if any(vals):
                texts.append("\t".join(vals))
    wb.close()
    return "\n".join(texts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


# ── チャンキング ───────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 1024, overlap: int = 200) -> list[str]:
    """固定長チャンク分割（フォールバック用）"""
    char_size = chunk_size * 2
    char_overlap = overlap * 2

    if len(text) <= char_size:
        return [text] if text.strip() else []

    chunks = []
    start = 0
    while start < len(text):
        end = start + char_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += char_size - char_overlap

    return chunks


import math
import re


def _split_sentences(text: str) -> list[str]:
    """テキストを文単位で分割（日本語・英語対応）"""
    # 段落（空行）→ 文末句読点で分割
    parts = re.split(r'(?<=[。！？\.\!\?\n])\s*', text)
    sentences = [s.strip() for s in parts if s.strip()]
    return sentences if sentences else [text.strip()] if text.strip() else []


def _cosine_sim(a: list[float], b: list[float]) -> float:
    """コサイン類似度（pure Python）"""
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(x * x for x in b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def semantic_chunk_text(
    text: str,
    max_chars: int = 2048,
    min_chars: int = 100,
    breakpoint_percentile: int = 25,
) -> list[str]:
    """セマンティックチャンキング: 文の埋め込み類似度で意味境界を検出して分割

    1. 文に分割
    2. 各文を埋め込み
    3. 隣接文の類似度が低い箇所（話題の変わり目）で分割
    4. max_chars / min_chars で調整
    """
    if not text.strip():
        return []
    if len(text) <= max_chars:
        return [text.strip()]

    sentences = _split_sentences(text)
    if len(sentences) <= 1:
        return [text.strip()]

    # 文のエンベディングを取得（バッチ）
    try:
        embeddings = get_embeddings(sentences)
    except Exception:
        log.warning("セマンティックチャンキング失敗 → 固定長フォールバック")
        return chunk_text(text)

    # 隣接文間のコサイン類似度
    similarities = []
    for i in range(len(embeddings) - 1):
        similarities.append(_cosine_sim(embeddings[i], embeddings[i + 1]))

    # ブレークポイント閾値（パーセンタイル）
    sorted_sims = sorted(similarities)
    idx = max(0, int(len(sorted_sims) * breakpoint_percentile / 100) - 1)
    threshold = sorted_sims[idx] if sorted_sims else 0.5

    # 意味境界で分割
    chunks = []
    current = sentences[0]

    for i in range(1, len(sentences)):
        candidate = current + "\n" + sentences[i]

        # 類似度が閾値未満 = 話題が変わった → 分割
        if similarities[i - 1] < threshold and len(current) >= min_chars:
            chunks.append(current.strip())
            current = sentences[i]
        # max_chars 超過 → 強制分割
        elif len(candidate) > max_chars and len(current) >= min_chars:
            chunks.append(current.strip())
            current = sentences[i]
        else:
            current = candidate

    if current.strip():
        chunks.append(current.strip())

    return [c for c in chunks if c]


# ── エンベディング ─────────────────────────────────────

@retry(
    retry=retry_if_exception(_is_transient_openai_error),
    wait=wait_exponential(multiplier=1, min=1, max=10),
    stop=stop_after_attempt(3),
    before_sleep=before_sleep_log(log, logging.WARNING),
    reraise=True,
)
def _embed_batch(client: AzureOpenAI, batch: list[str]) -> list[list[float]]:
    """エンベディングバッチ呼び出し（リトライ付き）"""
    resp = client.embeddings.create(input=batch, model=EMBEDDING_DEPLOYMENT)
    return [d.embedding for d in resp.data]


def get_embeddings(texts: list[str]) -> list[list[float]]:
    """Azure OpenAI でエンベディングを生成"""
    client = AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_KEY,
        api_version="2024-06-01",
    )
    # バッチサイズ制限（Azure OpenAI は 16 件/リクエスト推奨）
    all_embeddings = []
    batch_size = 16
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        all_embeddings.extend(_embed_batch(client, batch))
        if i + batch_size < len(texts):
            time.sleep(0.5)  # レート制限対策
    return all_embeddings


# ── DB 書き込み ────────────────────────────────────────

def upsert_chunks(
    file_id: str,
    title: str,
    source_url: str,
    category: str,
    allowed_groups: list[str],
    chunks: list[str],
    embeddings: list[list[float]],
):
    """チャンクを pgvector に upsert — トランザクション保証 (F-25)

    クラッシュ時は rollback → 旧チャンクが残る（データ消失なし）。
    """
    conn = get_conn()
    try:
        # プール返却時にトランザクション中の場合があるためリセット
        if conn.info.transaction_status != 0:  # IDLE 以外
            conn.rollback()
        conn.autocommit = False
        cur = conn.cursor()
        try:
            # 既存チャンクを削除
            cur.execute("DELETE FROM chunks WHERE file_id = %s", (file_id,))

            for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                chunk_id = hashlib.sha256(f"{file_id}:{i}".encode()).hexdigest()[:16]
                cur.execute(
                    """INSERT INTO chunks
                       (chunk_id, file_id, chunk_text, embedding, title, source_url, category, allowed_groups)
                       VALUES (%s, %s, %s, %s::vector, %s, %s, %s, %s)""",
                    (chunk_id, file_id, chunk, emb, title, source_url, category, allowed_groups),
                )

            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            cur.close()
    finally:
        conn.autocommit = True
        put_conn(conn)


# ── DB からの既存ファイル情報取得 ─────────────────────────

def _get_indexed_files() -> dict[str, str]:
    """DB にインジェスト済みのファイル ID と updated_at を取得"""
    conn = get_conn()
    try:
        cur = conn.cursor()
        cur.execute("SELECT DISTINCT file_id, updated_at FROM chunks")
        result = {row[0]: row[1].isoformat() if row[1] else "" for row in cur.fetchall()}
        cur.close()
        return result
    finally:
        put_conn(conn)


def _delete_file_chunks(file_id: str):
    """ファイルのチャンクを全削除（SP から消えたファイル用）"""
    conn = get_conn()
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM chunks WHERE file_id = %s", (file_id,))
        deleted = cur.rowcount
        conn.commit()
        cur.close()
        return deleted
    finally:
        put_conn(conn)


def _update_acl_only(file_id: str, allowed_groups: list[str]):
    """ACL だけ更新（ファイル内容は変わっていないが権限が変わった場合）"""
    conn = get_conn()
    try:
        cur = conn.cursor()
        cur.execute(
            "UPDATE chunks SET allowed_groups = %s, updated_at = now() WHERE file_id = %s",
            (allowed_groups, file_id),
        )
        conn.commit()
        cur.close()
    finally:
        put_conn(conn)


# ── メインパイプライン ─────────────────────────────────

_INGEST_LOCK_ID = 839201764  # アドバイザリーロック用の固有 ID


def run(incremental: bool = False):
    """インジェストパイプライン実行

    incremental=False: 全件再処理
    incremental=True:  差分更新（変更・追加・削除を検知）

    排他制御 (F-26): PostgreSQL アドバイザリーロックで二重実行を防止。
    """
    # ── 排他制御: アドバイザリーロック取得 ──
    lock_conn = get_conn()
    try:
        lock_cur = lock_conn.cursor()
        lock_cur.execute("SELECT pg_try_advisory_lock(%s)", (_INGEST_LOCK_ID,))
        acquired = lock_cur.fetchone()[0]
        lock_cur.close()
        if not acquired:
            log.warning("Another ingest process is running. Skipping.")
            put_conn(lock_conn)
            return
    except Exception:
        put_conn(lock_conn)
        raise

    try:
        _run_ingest(lock_conn, incremental)
    finally:
        # ロック解放 — 元の接続が切れている可能性があるため新規接続でも試行
        released = False
        for conn_attempt in [lock_conn, None]:
            try:
                c = conn_attempt if conn_attempt else get_conn()
                cur = c.cursor()
                cur.execute("SELECT pg_advisory_unlock(%s)", (_INGEST_LOCK_ID,))
                cur.close()
                if not conn_attempt:
                    put_conn(c)
                released = True
                break
            except Exception:
                if not conn_attempt:
                    log.exception("Failed to release advisory lock (retry also failed)")
                continue
        if not released:
            log.error("Advisory lock %s could not be released — manual cleanup required", _INGEST_LOCK_ID)
        try:
            put_conn(lock_conn)
        except Exception:
            pass


def _run_ingest(lock_conn, incremental: bool):
    """インジェスト本体（ロック取得後に実行）"""
    log.info("=== インジェスト開始 (incremental=%s) ===", incremental)
    token = get_app_token()
    _token_acquired_at = time.time()

    def _refresh_token_if_needed() -> str:
        """50分経過でトークンを自動更新（有効期限60分のマージン）"""
        nonlocal token, _token_acquired_at
        if time.time() - _token_acquired_at > 3000:  # 50分
            log.info("トークン更新（50分経過）")
            token = get_app_token()
            _token_acquired_at = time.time()
        return token

    # 権限チェック preflight — ACL 展開に必要な権限があるか確認
    check_graph_permissions(token)

    files = list_sp_files(token)
    log.info("SP ファイル数: %d", len(files))

    sp_file_ids = {f["id"] for f in files}

    # 差分検知用: DB の既存ファイル
    indexed = _get_indexed_files() if incremental else {}

    processed = 0
    skipped = 0
    acl_updated = 0
    deleted = 0
    skip_reasons: list[tuple[str, str]] = []  # (filename, reason)

    for f in files:
        token = _refresh_token_if_needed()
        title = f["name"]
        path = f["path"]
        category = path.split("/")[0] if path else "root"
        source_url = f["webUrl"]
        file_id = f["id"]
        last_modified = f["lastModified"]

        # 差分チェック: incremental モードでは変更がないファイルをスキップ
        if incremental and file_id in indexed:
            # ACL だけ更新チェック
            current_acl = resolve_folder_acl(token, path)
            _update_acl_only(file_id, current_acl)
            acl_updated += 1

            # ファイル内容が変わっていなければスキップ
            # lastModifiedDateTime と DB の updated_at を比較
            if last_modified <= indexed[file_id]:
                skipped += 1
                continue

        log.info("処理中: %s/%s", path, title)

        # テキスト抽出
        content = download_file(token, file_id)
        text = extract_text(content, title)
        if not text.strip():
            log.warning("  テキスト抽出結果が空: %s", title)
            ext = title.rsplit(".", 1)[-1].lower() if "." in title else ""
            if ext in ("doc", "xls", "ppt"):
                skip_reasons.append((title, f"旧Office形式 (.{ext})"))
            elif ext == "pdf":
                skip_reasons.append((title, "テキスト抽出不可（画像のみPDF/パスワード保護の可能性）"))
            else:
                skip_reasons.append((title, "テキスト空"))
            skipped += 1
            continue

        # ACL 取得
        allowed_groups = resolve_folder_acl(token, path)
        log.info("  ACL: %s", allowed_groups)

        # セマンティックチャンキング
        chunks = semantic_chunk_text(text)
        # チャンク品質フィルタ: 意味のある文字が少なすぎるチャンクを除外
        before_filter = len(chunks)
        chunks = [c for c in chunks if len(re.sub(r'[\.\s\d\u3000]+', '', c)) >= 10]
        if before_filter != len(chunks):
            log.info("  品質フィルタ: %d → %d チャンク", before_filter, len(chunks))
        log.info("  チャンク数: %d (semantic)", len(chunks))

        if not chunks:
            skipped += 1
            continue

        # エンベディング
        embeddings = get_embeddings(chunks)

        # DB 書き込み
        upsert_chunks(file_id, title, source_url, category, allowed_groups, chunks, embeddings)
        processed += 1

    # SP から削除されたファイルを DB からも削除
    if incremental and indexed:
        for old_id in indexed:
            if old_id not in sp_file_ids:
                count = _delete_file_chunks(old_id)
                log.info("削除: file_id=%s (%d chunks)", old_id, count)
                deleted += count

    # フルインジェスト時: TARGET_FOLDERS に含まれないカテゴリのチャンクを削除
    if not incremental and TARGET_FOLDERS:
        conn = get_conn()
        try:
            cur = conn.cursor()
            cur.execute("SELECT DISTINCT category FROM chunks")
            all_cats = [row[0] for row in cur.fetchall()]
            for cat in all_cats:
                if not any(cat.startswith(prefix) for prefix in TARGET_FOLDERS):
                    cur.execute("DELETE FROM chunks WHERE category = %s", (cat,))
                    orphan_count = cur.rowcount
                    if orphan_count:
                        log.info("対象外カテゴリ削除: %s (%d chunks)", cat, orphan_count)
                        deleted += orphan_count
            conn.commit()
            cur.close()
        finally:
            put_conn(conn)

    log.info(
        "=== インジェスト完了: %d 処理, %d スキップ, %d ACL更新, %d チャンク削除 ===",
        processed, skipped, acl_updated, deleted,
    )

    # スキップレポート
    if skip_reasons:
        log.info("=== スキップレポート (%d 件) ===", len(skip_reasons))
        for fname, reason in skip_reasons:
            log.info("  %s: %s", fname, reason)


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
    incremental = "--incremental" in sys.argv
    run(incremental=incremental)
